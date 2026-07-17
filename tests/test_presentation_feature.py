from __future__ import annotations

import json
import tempfile
import unittest

from pathlib import Path
from unittest.mock import patch

from pptx import Presentation
from pptx.util import Inches
from pydantic import ValidationError

from reins.features.presentation.engines.frontend_slides import (
    FrontendSlidesEngine,
)
from reins.features.presentation.models import (
    PresentationEngine,
    PresentationJobStatus,
    PresentationOutputFormat,
    PresentationRequest,
    PresentationStyle,
)
from reins.features.presentation.planner import create_basic_plan
from reins.features.presentation.router import select_presentation_engine
from reins.features.presentation.service import PresentationService
from reins.features.presentation.storage import PresentationStorage


def presentation_request(**overrides) -> PresentationRequest:
    payload = {
        "prompt": (
            "Create a product launch deck for a privacy-first AI assistant. "
            "Cover the customer problem, solution, differentiation, rollout, "
            "risks, and decisions required."
        ),
        "title": "Privacy-first AI launch",
        "audience": "Product leadership",
        "slide_count": 6,
        "style": "modern",
        "engine": "native_pptx",
        "output_format": "pptx",
        "metadata": {"skip_ai": True},
    }
    payload.update(overrides)
    return PresentationRequest.model_validate(payload)


def create_source_deck(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.6), Inches(8), Inches(0.8)
    )
    title.text = "Quarterly Review"
    body = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.8), Inches(8), Inches(2)
    )
    body.text = "Old metric\nKeep this sentence"
    presentation.core_properties.title = "Quarterly Review"
    presentation.save(path)


class PresentationFeatureTests(unittest.TestCase):
    def test_new_presentation_requires_a_prompt(self):
        with self.assertRaises(ValidationError):
            PresentationRequest(prompt="   ")

    def test_router_selects_frontend_slides_for_html(self):
        request = presentation_request(
            engine="auto",
            output_format="html",
        )

        self.assertEqual(
            select_presentation_engine(request),
            PresentationEngine.FRONTEND_SLIDES,
        )

    def test_basic_planner_creates_requested_slide_count(self):
        plan = create_basic_plan(presentation_request())

        self.assertEqual(plan.title, "Privacy-first AI launch")
        self.assertEqual(len(plan.slides), 6)
        self.assertEqual(plan.slides[0].type, "title")
        self.assertEqual(plan.slides[-1].type, "conclusion")
        self.assertTrue(all(slide.title for slide in plan.slides))

    def test_native_service_creates_valid_pptx(self):
        with tempfile.TemporaryDirectory() as directory:
            service = PresentationService(
                PresentationStorage(Path(directory))
            )

            result = service.create_job(presentation_request())

            self.assertTrue(result.success)
            self.assertIsNotNone(result.primary_output)
            assert result.primary_output is not None
            self.assertTrue(result.primary_output.is_file())
            self.assertEqual(len(Presentation(result.primary_output).slides), 6)

            state = service.get_job_state(result.job_id or "")
            self.assertEqual(state.status, PresentationJobStatus.COMPLETED)
            self.assertEqual(state.progress, 100)
            self.assertEqual(state.output_path, result.primary_output)
            self.assertTrue(
                any(artifact.kind == "qa-report" for artifact in state.artifacts)
            )

            assert state.output_path is not None
            report = json.loads(
                (state.output_path.parent.parent / "qa-report.json").read_text(
                    encoding="utf-8"
                )
            )
            self.assertTrue(report["ok"])

    def test_auto_uses_native_renderer_for_four_by_three(self):
        with tempfile.TemporaryDirectory() as directory:
            service = PresentationService(
                PresentationStorage(Path(directory))
            )
            request = presentation_request(
                engine="auto",
                aspect_ratio="4:3",
            )

            result = service.create_job(request)

            self.assertTrue(result.success)
            self.assertEqual(result.engine, PresentationEngine.NATIVE_PPTX)
            self.assertTrue(any("4:3" in warning for warning in result.warnings))

    def test_new_presentation_uses_uploaded_pdf_as_source_material(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "research.pdf"
            source.write_bytes(b"%PDF-test")
            service = PresentationService(PresentationStorage(root / "jobs"))

            def extract_source(_source: Path, output: Path) -> str:
                content = "Revenue grew 42 percent. Retention reached 91 percent."
                output.write_text(content, encoding="utf-8")
                return content

            with patch(
                "reins.features.presentation.service.extract_pdf_markdown",
                side_effect=extract_source,
            ):
                result = service.create_job(
                    presentation_request(
                        source_path=source,
                        prompt="Build a board review from the uploaded research.",
                    )
                )

            self.assertTrue(result.success, result.errors)
            plan = service.get_job_plan(result.job_id or "")
            self.assertIn("Revenue grew 42 percent", plan.metadata["source_prompt"])
            self.assertTrue(
                any(artifact.kind == "source-markdown" for artifact in result.artifacts)
            )

    def test_frontend_slides_escapes_user_content(self):
        with tempfile.TemporaryDirectory() as directory:
            engine = FrontendSlidesEngine()
            request = presentation_request(
                title="Safety <script>alert(1)</script>",
                engine="frontend_slides",
                output_format="html",
                style="creative",
            )
            plan = create_basic_plan(request)

            result = engine.render(
                request=request,
                plan=plan,
                workspace=Path(directory),
            )

            self.assertTrue(result.success)
            self.assertIsNotNone(result.primary_output)
            assert result.primary_output is not None
            source = result.primary_output.read_text(encoding="utf-8")
            self.assertIn(
                "Safety &lt;script&gt;alert(1)&lt;/script&gt;",
                source,
            )
            self.assertNotIn("Safety <script>alert(1)</script>", source)
            self.assertEqual(source.count('<section class="slide '), 6)
            self.assertIn("ArrowRight", source)

    def test_output_format_enum_remains_limited(self):
        self.assertEqual(
            set(PresentationOutputFormat),
            {
                PresentationOutputFormat.PPTX,
                PresentationOutputFormat.HTML,
                PresentationOutputFormat.PDF,
            },
        )
        self.assertEqual(PresentationStyle.TECH.value, "tech")

    def test_modify_applies_exact_text_change_and_preserves_other_parts(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "source.pptx"
            create_source_deck(source)
            service = PresentationService(PresentationStorage(root / "jobs"))
            request = presentation_request(
                action="modify",
                prompt=None,
                source_path=source,
                instruction='replace "Old metric" with "Updated metric"',
                slide_count=1,
                metadata={"skip_ai": True},
            )

            result = service.create_job(request)

            self.assertTrue(result.success, result.errors)
            assert result.primary_output is not None
            edited = Presentation(result.primary_output)
            slide_text = "\n".join(
                shape.text for shape in edited.slides[0].shapes if shape.has_text_frame
            )
            self.assertIn("Updated metric", slide_text)
            self.assertNotIn("Old metric", slide_text)
            report_path = next(
                artifact.path
                for artifact in result.artifacts
                if artifact.kind == "preservation-report"
            )
            report = json.loads(report_path.read_text(encoding="utf-8"))
            self.assertTrue(report["ok"])
            self.assertEqual(report["changed_parts"], ["ppt/slides/slide1.xml"])

    def test_restyle_changes_theme_without_rewriting_slides(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "source.pptx"
            create_source_deck(source)
            service = PresentationService(PresentationStorage(root / "jobs"))
            request = presentation_request(
                action="restyle",
                prompt=None,
                source_path=source,
                instruction="Use a restrained dark technology theme.",
                style="dark",
                slide_count=1,
            )

            result = service.create_job(request)

            self.assertTrue(result.success, result.errors)
            report_path = next(
                artifact.path
                for artifact in result.artifacts
                if artifact.kind == "preservation-report"
            )
            report = json.loads(report_path.read_text(encoding="utf-8"))
            self.assertTrue(report["ok"])
            self.assertTrue(report["changed_parts"])
            self.assertTrue(
                all(part.startswith("ppt/theme/theme") for part in report["changed_parts"])
            )

    def test_convert_existing_pptx_to_semantic_html(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "source.pptx"
            create_source_deck(source)
            service = PresentationService(PresentationStorage(root / "jobs"))
            request = presentation_request(
                action="convert",
                prompt=None,
                source_path=source,
                instruction="Create an accessible web version.",
                output_format="html",
                engine="auto",
                slide_count=1,
            )

            result = service.create_job(request)

            self.assertTrue(result.success, result.errors)
            assert result.primary_output is not None
            html = result.primary_output.read_text(encoding="utf-8")
            self.assertIn("Quarterly Review", html)
            self.assertIn("Old metric", html)
            self.assertIn("ArrowRight", html)

    def test_modify_restyle_convert_revision_chain(self):
        with tempfile.TemporaryDirectory() as directory:
            root = Path(directory)
            source = root / "source.pptx"
            create_source_deck(source)
            service = PresentationService(PresentationStorage(root / "jobs"))

            modified = service.create_job(
                presentation_request(
                    action="modify",
                    prompt=None,
                    source_path=source,
                    instruction='replace "Old metric" with "Final metric"',
                    slide_count=1,
                    metadata={"skip_ai": True},
                )
            )
            self.assertTrue(modified.success, modified.errors)
            assert modified.primary_output is not None

            restyled = service.create_job(
                presentation_request(
                    action="restyle",
                    prompt=None,
                    source_path=modified.primary_output,
                    instruction="Use a dark technology theme.",
                    style="modern",
                    slide_count=1,
                )
            )
            self.assertTrue(restyled.success, restyled.errors)
            assert restyled.primary_output is not None

            converted = service.create_job(
                presentation_request(
                    action="convert",
                    prompt=None,
                    source_path=restyled.primary_output,
                    instruction="Create an accessible web version.",
                    output_format="html",
                    engine="auto",
                    slide_count=1,
                )
            )
            self.assertTrue(converted.success, converted.errors)
            assert converted.primary_output is not None
            html = converted.primary_output.read_text(encoding="utf-8")
            self.assertIn("Final metric", html)

            restyle_report_path = next(
                artifact.path
                for artifact in restyled.artifacts
                if artifact.kind == "preservation-report"
            )
            restyle_report = json.loads(
                restyle_report_path.read_text(encoding="utf-8")
            )
            self.assertEqual(restyle_report["style"], "dark")


if __name__ == "__main__":
    unittest.main()
