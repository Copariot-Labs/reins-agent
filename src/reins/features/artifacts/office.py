from __future__ import annotations

import json
import re
from pathlib import Path
from typing import Any

from reins.features.artifacts.paths import unique_artifact_path
from reins.features.artifacts.schema import ArtifactRecord, normalize_artifact_format, normalize_title
from reins.features.artifacts.store import ArtifactStore, get_default_artifact_store


_XML_INVALID_RE = re.compile(
    r"[\x00-\x08\x0B\x0C\x0E-\x1F\uD800-\uDFFF\uFFFE\uFFFF]"
)

_MARKDOWN_BOLD_RE = re.compile(r"\*\*(.*?)\*\*")


def sanitize_xml_text(value: object) -> str:
    text = str(value or "")
    text = text.replace("\\x00", "")
    text = text.replace("\x00", "")
    text = _XML_INVALID_RE.sub("", text)
    return text


def clean_inline_markdown(value: object) -> str:
    text = sanitize_xml_text(value)
    text = text.replace("**", "")
    text = text.replace("__", "")
    text = text.replace("`", "")
    text = re.sub(r"^\s*#+\s*", "", text)
    return text.strip()


def _normalize_paragraphs(body: str | list[str]) -> list[str]:
    if isinstance(body, list):
        paragraphs: list[str] = []

        for item in body:
            cleaned = sanitize_xml_text(item).strip()

            if cleaned:
                paragraphs.append(cleaned)

        return paragraphs

    text = sanitize_xml_text(body).strip()

    if not text:
        return []

    return [part.strip() for part in re.split(r"\n\s*\n", text) if part.strip()]


def _ensure_parent(path: str | Path) -> Path:
    output_path = Path(path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    return output_path


def _safe_sheet_name(value: object, *, fallback: str = "Sheet") -> str:
    name = sanitize_xml_text(value).strip() or fallback
    name = re.sub(r"[\[\]\:\*\?\/\\]", "-", name)
    name = name[:31].strip()
    return name or fallback


def _parse_markdown_table(lines: list[str]) -> tuple[list[str], list[list[str]]] | None:
    table_lines = [
        line.strip()
        for line in lines
        if line.strip().startswith("|") and line.strip().endswith("|")
    ]

    if len(table_lines) < 2:
        return None

    rows: list[list[str]] = []

    for line in table_lines:
        cells = [clean_inline_markdown(cell.strip()) for cell in line.strip("|").split("|")]

        if all(re.fullmatch(r":?-{3,}:?", cell.replace(" ", "")) for cell in cells):
            continue

        rows.append(cells)

    if not rows:
        return None

    headers = rows[0]
    body_rows = rows[1:]

    return headers, body_rows


def _looks_like_heading(line: str) -> tuple[int, str] | None:
    cleaned = sanitize_xml_text(line).strip()

    if not cleaned:
        return None

    markdown_match = re.match(r"^(#{1,4})\s+(.+)$", cleaned)
    if markdown_match:
        level = min(len(markdown_match.group(1)), 4)
        return level, clean_inline_markdown(markdown_match.group(2))

    roman_match = re.match(r"^([IVXLCDM]+)\.\s+(.+)$", cleaned)
    if roman_match:
        return 2, clean_inline_markdown(cleaned)

    numbered_section_match = re.match(r"^\d+(\.\d+)*\s+.+$", cleaned)
    if numbered_section_match and len(cleaned) <= 100:
        return 2, clean_inline_markdown(cleaned)

    if cleaned.endswith(":") and len(cleaned) <= 80:
        return 3, clean_inline_markdown(cleaned[:-1])

    return None


def _looks_like_bullet(line: str) -> str | None:
    match = re.match(r"^\s*[-*•]\s+(.+)$", line)

    if not match:
        return None

    return clean_inline_markdown(match.group(1))


def _looks_like_numbered_item(line: str) -> str | None:
    match = re.match(r"^\s*\d+[\.\)]\s+(.+)$", line)

    if not match:
        return None

    return clean_inline_markdown(match.group(1))


def _add_inline_runs(paragraph, text: str) -> None:
    cleaned = sanitize_xml_text(text)

    if not cleaned:
        return

    position = 0

    for match in _MARKDOWN_BOLD_RE.finditer(cleaned):
        before = cleaned[position : match.start()]
        bold_text = match.group(1)

        if before:
            paragraph.add_run(before.replace("`", ""))

        run = paragraph.add_run(sanitize_xml_text(bold_text).replace("`", ""))
        run.bold = True
        position = match.end()

    remaining = cleaned[position:]

    if remaining:
        paragraph.add_run(
            remaining.replace("**", "").replace("__", "").replace("`", "")
        )


def _apply_document_defaults(document) -> None:
    from docx.shared import Pt

    styles = document.styles

    normal = styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)

    for style_name in ["Heading 1", "Heading 2", "Heading 3"]:
        if style_name in styles:
            styles[style_name].font.name = "Calibri"


def _add_docx_table(document, lines: list[str]) -> bool:
    parsed = _parse_markdown_table(lines)

    if not parsed:
        return False

    headers, rows = parsed

    if not headers:
        return False

    table = document.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"

    header_cells = table.rows[0].cells

    for index, header in enumerate(headers):
        if index >= len(header_cells):
            break

        paragraph = header_cells[index].paragraphs[0]
        run = paragraph.add_run(clean_inline_markdown(header))
        run.bold = True

    for row in rows:
        cells = table.add_row().cells

        for index, value in enumerate(row[: len(cells)]):
            cells[index].text = clean_inline_markdown(value)

    document.add_paragraph()
    return True


def create_docx_file(
    *,
    title: str,
    body: str | list[str],
    output_path: str | Path,
) -> Path:
    try:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Inches, Pt
    except Exception as exc:
        raise RuntimeError(
            "python-docx is required to create DOCX artifacts. "
            "Install it with `pip install python-docx`."
        ) from exc

    path = _ensure_parent(output_path)

    document = Document()
    _apply_document_defaults(document)

    section = document.sections[0]
    section.top_margin = Inches(0.75)
    section.bottom_margin = Inches(0.75)
    section.left_margin = Inches(0.85)
    section.right_margin = Inches(0.85)

    title_text = clean_inline_markdown(
        normalize_title(title, default="Untitled Document")
    )

    title_paragraph = document.add_paragraph()
    title_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER

    title_run = title_paragraph.add_run(title_text)
    title_run.bold = True
    title_run.font.size = Pt(18)

    document.add_paragraph()

    raw_text = "\n\n".join(_normalize_paragraphs(body))
    lines = [sanitize_xml_text(line).rstrip() for line in raw_text.splitlines()]

    buffer: list[str] = []
    index = 0

    def flush_buffer() -> None:
        nonlocal buffer

        if not buffer:
            return

        paragraph_text = " ".join(
            clean_inline_markdown(line) for line in buffer if line.strip()
        ).strip()

        if paragraph_text:
            paragraph = document.add_paragraph()
            paragraph.paragraph_format.space_after = Pt(6)
            paragraph.paragraph_format.line_spacing = 1.08
            _add_inline_runs(paragraph, paragraph_text)

        buffer = []

    while index < len(lines):
        line = lines[index].strip()

        if not line:
            flush_buffer()
            index += 1
            continue

        table_block: list[str] = []

        while (
            index < len(lines)
            and lines[index].strip().startswith("|")
            and lines[index].strip().endswith("|")
        ):
            table_block.append(lines[index])
            index += 1

        if table_block:
            flush_buffer()

            if not _add_docx_table(document, table_block):
                for table_line in table_block:
                    paragraph = document.add_paragraph(clean_inline_markdown(table_line))
                    paragraph.paragraph_format.space_after = Pt(4)

            continue

        heading = _looks_like_heading(line)
        if heading:
            flush_buffer()
            level, heading_text = heading
            heading_paragraph = document.add_heading(heading_text, level=min(level, 3))
            heading_paragraph.paragraph_format.space_before = Pt(8)
            heading_paragraph.paragraph_format.space_after = Pt(4)
            index += 1
            continue

        bullet = _looks_like_bullet(line)
        if bullet:
            flush_buffer()
            paragraph = document.add_paragraph(style="List Bullet")
            _add_inline_runs(paragraph, bullet)
            index += 1
            continue

        numbered = _looks_like_numbered_item(line)
        if numbered:
            flush_buffer()
            paragraph = document.add_paragraph(style="List Number")
            _add_inline_runs(paragraph, numbered)
            index += 1
            continue

        horizontal_rule = re.fullmatch(r"[-*_]{3,}", line)
        if horizontal_rule:
            flush_buffer()
            document.add_paragraph()
            index += 1
            continue

        buffer.append(line)
        index += 1

    flush_buffer()

    document.save(path)
    return path


def create_xlsx_file(
    *,
    title: str,
    sheets: list[dict[str, Any]] | None,
    output_path: str | Path,
) -> Path:
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font, PatternFill
    except Exception as exc:
        raise RuntimeError(
            "openpyxl is required to create XLSX artifacts. "
            "Install it with `pip install openpyxl`."
        ) from exc

    path = _ensure_parent(output_path)

    workbook = Workbook()

    normalized_sheets = sheets or [
        {
            "name": "Sheet1",
            "headers": ["Title", "Value"],
            "rows": [[normalize_title(title, default="Untitled Spreadsheet"), ""]],
        }
    ]

    first_sheet = True

    for sheet_index, sheet in enumerate(normalized_sheets, start=1):
        name = _safe_sheet_name(
            sheet.get("name") or f"Sheet{sheet_index}",
            fallback=f"Sheet{sheet_index}",
        )

        if first_sheet:
            worksheet = workbook.active
            worksheet.title = name
            first_sheet = False
        else:
            worksheet = workbook.create_sheet(title=name)

        headers = sheet.get("headers") or []
        rows = sheet.get("rows") or []

        if headers:
            worksheet.append([sanitize_xml_text(header) for header in headers])

            for cell in worksheet[1]:
                cell.font = Font(bold=True)
                cell.alignment = Alignment(
                    horizontal="center",
                    vertical="center",
                    wrap_text=True,
                )
                cell.fill = PatternFill(fill_type="solid", fgColor="D9EAF7")

        for row in rows:
            if isinstance(row, dict):
                if headers:
                    worksheet.append(
                        [sanitize_xml_text(row.get(header, "")) for header in headers]
                    )
                else:
                    worksheet.append(
                        [sanitize_xml_text(value) for value in row.values()]
                    )
            elif isinstance(row, list | tuple):
                worksheet.append([sanitize_xml_text(cell) for cell in row])
            else:
                worksheet.append([sanitize_xml_text(row)])

        worksheet.freeze_panes = "A2"

        for row in worksheet.iter_rows():
            for cell in row:
                cell.alignment = Alignment(vertical="top", wrap_text=True)

        for column_cells in worksheet.columns:
            max_length = 0
            column_letter = column_cells[0].column_letter

            for cell in column_cells:
                value = sanitize_xml_text(cell.value)
                max_length = max(max_length, len(value))

            worksheet.column_dimensions[column_letter].width = min(
                max(max_length + 2, 12),
                55,
            )

    workbook.save(path)
    return path


def create_pptx_file(
    *,
    title: str,
    slides: list[dict[str, Any]] | None,
    output_path: str | Path,
) -> Path:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise RuntimeError(
            "python-pptx is required to create PPTX artifacts. "
            "Install it with `pip install python-pptx`."
        ) from exc

    path = _ensure_parent(output_path)

    presentation = Presentation()
    normalized_title = sanitize_xml_text(
        normalize_title(title, default="Untitled Presentation")
    )

    normalized_slides = slides or [
        {
            "title": normalized_title,
            "bullets": ["Generated by Reins."],
        }
    ]

    for index, slide_data in enumerate(normalized_slides):
        slide_title = sanitize_xml_text(
            slide_data.get("title") or normalized_title
        ).strip() or normalized_title

        bullets = slide_data.get("bullets") or []
        body = sanitize_xml_text(slide_data.get("body") or "").strip()

        if index == 0:
            layout = presentation.slide_layouts[0]
            slide = presentation.slides.add_slide(layout)
            slide.shapes.title.text = slide_title

            subtitle = slide.placeholders[1]

            if bullets:
                subtitle.text = "\n".join(
                    sanitize_xml_text(item) for item in bullets[:5]
                )
            elif body:
                subtitle.text = body
            else:
                subtitle.text = ""

            continue

        layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = slide_title

        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()

        normalized_bullets = bullets or ([body] if body else [""])

        for bullet_index, bullet in enumerate(normalized_bullets):
            paragraph = (
                text_frame.paragraphs[0]
                if bullet_index == 0
                else text_frame.add_paragraph()
            )
            paragraph.text = sanitize_xml_text(bullet)
            paragraph.level = 0

    presentation.save(path)
    return path


def create_txt_file(
    *,
    title: str,
    body: str | list[str],
    output_path: str | Path,
) -> Path:
    path = _ensure_parent(output_path)

    content = "\n\n".join(
        [
            clean_inline_markdown(normalize_title(title, default="Untitled Text")),
            "\n\n".join(_normalize_paragraphs(body)),
        ]
    ).strip()

    path.write_text(content + "\n", encoding="utf-8")
    return path


def _sanitize_json_payload(value: Any) -> Any:
    if isinstance(value, dict):
        return {
            sanitize_xml_text(key): _sanitize_json_payload(item)
            for key, item in value.items()
        }

    if isinstance(value, list):
        return [_sanitize_json_payload(item) for item in value]

    if isinstance(value, tuple):
        return [_sanitize_json_payload(item) for item in value]

    if isinstance(value, str):
        return sanitize_xml_text(value)

    return value


def create_json_file(
    *,
    payload: dict[str, Any],
    output_path: str | Path,
) -> Path:
    path = _ensure_parent(output_path)
    safe_payload = _sanitize_json_payload(payload)

    path.write_text(
        json.dumps(safe_payload, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )
    return path


def create_office_artifact(
    *,
    title: str,
    body: str | list[str] = "",
    artifact_format: str = "docx",
    sheets: list[dict[str, Any]] | None = None,
    slides: list[dict[str, Any]] | None = None,
    metadata: dict[str, Any] | None = None,
    store: ArtifactStore | None = None,
    source: str = "reins",
) -> ArtifactRecord:
    artifact_store = store or get_default_artifact_store()
    normalized_format = normalize_artifact_format(artifact_format, default="docx")
    normalized_title = clean_inline_markdown(normalize_title(title))

    output_path = unique_artifact_path(
        title=normalized_title,
        artifact_format=normalized_format,
    )

    safe_body: str | list[str]

    if isinstance(body, list):
        safe_body = [sanitize_xml_text(item) for item in body]
    else:
        safe_body = sanitize_xml_text(body)

    safe_sheets = _sanitize_json_payload(sheets or [])
    safe_slides = _sanitize_json_payload(slides or [])
    safe_metadata = _sanitize_json_payload(metadata or {})

    if normalized_format == "docx":
        path = create_docx_file(
            title=normalized_title,
            body=safe_body,
            output_path=output_path,
        )
    elif normalized_format == "xlsx":
        path = create_xlsx_file(
            title=normalized_title,
            sheets=safe_sheets,
            output_path=output_path,
        )
    elif normalized_format == "pptx":
        path = create_pptx_file(
            title=normalized_title,
            slides=safe_slides,
            output_path=output_path,
        )
    elif normalized_format == "json":
        path = create_json_file(
            payload={
                "title": normalized_title,
                "body": safe_body,
                "sheets": safe_sheets,
                "slides": safe_slides,
                "metadata": safe_metadata,
            },
            output_path=output_path,
        )
    else:
        path = create_txt_file(
            title=normalized_title,
            body=safe_body,
            output_path=output_path,
        )

    return artifact_store.register(
        title=normalized_title,
        kind=normalized_format,
        path=path,
        summary=f"Created {normalized_format.upper()} artifact: {Path(path).name}",
        source=source,
        metadata={
            "artifact_format": normalized_format,
            "body": safe_body,
            "sheets": safe_sheets,
            "slides": safe_slides,
            **dict(safe_metadata or {}),
        },
    )


def create_docx_artifact(
    *,
    title: str,
    body: str | list[str],
    metadata: dict[str, Any] | None = None,
    store: ArtifactStore | None = None,
    source: str = "reins",
) -> ArtifactRecord:
    return create_office_artifact(
        title=title,
        body=body,
        artifact_format="docx",
        metadata=metadata,
        store=store,
        source=source,
    )


def create_xlsx_artifact(
    *,
    title: str,
    sheets: list[dict[str, Any]],
    metadata: dict[str, Any] | None = None,
    store: ArtifactStore | None = None,
    source: str = "reins",
) -> ArtifactRecord:
    return create_office_artifact(
        title=title,
        artifact_format="xlsx",
        sheets=sheets,
        metadata=metadata,
        store=store,
        source=source,
    )


def create_pptx_artifact(
    *,
    title: str,
    slides: list[dict[str, Any]],
    metadata: dict[str, Any] | None = None,
    store: ArtifactStore | None = None,
    source: str = "reins",
) -> ArtifactRecord:
    return create_office_artifact(
        title=title,
        artifact_format="pptx",
        slides=slides,
        metadata=metadata,
        store=store,
        source=source,
    )