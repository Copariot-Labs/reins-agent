from __future__ import annotations

from datetime import datetime
import importlib.util
from pathlib import Path
import re
import shutil
from typing import Any

from reins.api.home import get_reins_home


class ArtifactDependencyError(RuntimeError):
    pass


def check_artifact_dependencies() -> dict[str, bool]:
    playwright_browser = False

    if importlib.util.find_spec("playwright") is not None:
        try:
            from playwright.sync_api import sync_playwright

            with sync_playwright() as p:
                playwright_browser = Path(p.chromium.executable_path).exists()
        except Exception:
            playwright_browser = False

    return {
        "python-docx": importlib.util.find_spec("docx") is not None,
        "openpyxl": importlib.util.find_spec("openpyxl") is not None,
        "python-pptx": importlib.util.find_spec("pptx") is not None,
        "playwright": importlib.util.find_spec("playwright") is not None,
        "playwright-chromium": playwright_browser,
        "pillow": importlib.util.find_spec("PIL") is not None,
        "pytesseract": importlib.util.find_spec("pytesseract") is not None,
        "tesseract": shutil.which("tesseract") is not None,
    }


def _load_document_class():
    try:
        from docx import Document
    except ImportError as exc:
        raise ArtifactDependencyError(
            "Missing dependency: python-docx. Install Reins with project dependencies "
            "before running Office workmode tasks."
        ) from exc

    return Document


def _load_workbook_class():
    try:
        from openpyxl import Workbook
    except ImportError as exc:
        raise ArtifactDependencyError(
            "Missing dependency: openpyxl. Install Reins with project dependencies "
            "before running Excel workmode tasks."
        ) from exc

    return Workbook


def _load_presentation_class():
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ArtifactDependencyError(
            "Missing dependency: python-pptx. Install Reins with project dependencies "
            "before running PowerPoint workmode tasks."
        ) from exc

    return Presentation


def get_artifact_dir() -> Path:
    path = get_reins_home() / "workmode" / "artifacts"
    path.mkdir(parents=True, exist_ok=True)
    return path


SUPPORTED_OFFICE_FORMATS = {"docx", "xlsx", "pptx"}


def normalize_office_format(value: Any) -> str | None:
    text = str(value or "").strip().lower().lstrip(".")

    aliases = {
        "word": "docx",
        "document": "docx",
        "doc": "docx",
        "docx": "docx",
        "excel": "xlsx",
        "spreadsheet": "xlsx",
        "workbook": "xlsx",
        "table": "xlsx",
        "ledger": "xlsx",
        "xlsx": "xlsx",
        "xls": "xlsx",
        "powerpoint": "pptx",
        "presentation": "pptx",
        "slides": "pptx",
        "slide": "pptx",
        "deck": "pptx",
        "ppt": "pptx",
        "pptx": "pptx",
    }

    return aliases.get(text)


def infer_office_artifact_format(message: str = "", *hints: Any) -> str:
    for hint in hints:
        if isinstance(hint, dict):
            for key in (
                "artifact_format",
                "file_type",
                "format",
                "artifact_kind",
                "document_format",
            ):
                normalized = normalize_office_format(hint.get(key))
                if normalized:
                    return normalized
            for key in ("expected_artifacts", "artifacts"):
                values = hint.get(key)
                if isinstance(values, list):
                    for value in values:
                        normalized = normalize_office_format(value)
                        if normalized:
                            return normalized
            continue

        if isinstance(hint, (list, tuple, set)):
            for value in hint:
                normalized = normalize_office_format(value)
                if normalized:
                    return normalized
            continue

        normalized = normalize_office_format(hint)
        if normalized:
            return normalized

    text = str(message or "").lower()

    if any(token in text for token in ("台账", "表格")) or re.search(
        r"\b(xlsx|xls|excel|spreadsheet|workbook|ledger|table)\b",
        text,
    ):
        return "xlsx"

    if any(token in text for token in ("演示", "幻灯片")) or re.search(
        r"\b(pptx|ppt|powerpoint|presentation|slides?|slide deck|deck)\b",
        text,
    ):
        return "pptx"

    return "docx"


def _safe_title_slug(title: str) -> str:
    slug = re.sub(r"[^A-Za-z0-9._-]+", "-", str(title or "").strip()).strip("-._")
    return (slug[:48] or "office-artifact").lower()


def _artifact_path(title: str, suffix: str) -> Path:
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    return get_artifact_dir() / f"{_safe_title_slug(title)}-{timestamp}.{suffix}"


def _body_lines(body: str) -> list[str]:
    return [line.strip() for line in str(body or "").splitlines() if line.strip()]


def _body_lines_without_duplicate_title(title: str, body: str) -> list[str]:
    lines = _body_lines(body)
    if lines and lines[0].strip().lower() == str(title or "").strip().lower():
        return lines[1:]
    return lines


def generate_docx_artifact(title: str, body: str, content: dict[str, Any] | None = None) -> Path:
    path = _artifact_path(title or "document", "docx")
    Document = _load_document_class()

    doc = Document()
    doc.add_heading(title or "WorkMode Document", level=1)

    for line in _body_lines_without_duplicate_title(title, body):
        if line.startswith(("- ", "* ")):
            doc.add_paragraph(line[2:].strip(), style="List Bullet")
        elif re.match(r"^\d+[.)]\s+", line):
            doc.add_paragraph(re.sub(r"^\d+[.)]\s+", "", line).strip(), style="List Number")
        else:
            doc.add_paragraph(line)

    missing_fields = (content or {}).get("missing_fields")
    if isinstance(missing_fields, list) and missing_fields:
        doc.add_heading("Missing Fields", level=2)
        for field in missing_fields:
            doc.add_paragraph(str(field), style="List Bullet")

    doc.add_paragraph("Generated by Reins WorkMode.")
    doc.save(path)

    return path


def _normalize_sheet_rows(sheet: dict[str, Any], body: str) -> tuple[list[str], list[list[Any]]]:
    columns = sheet.get("columns")
    rows = sheet.get("rows")

    if isinstance(rows, list) and rows:
        if isinstance(columns, list) and columns:
            headers = [str(column) for column in columns]
        elif isinstance(rows[0], dict):
            headers = [str(key) for key in rows[0].keys()]
        else:
            width = max(len(row) if isinstance(row, list) else 1 for row in rows)
            headers = [f"Column {index}" for index in range(1, width + 1)]

        normalized_rows: list[list[Any]] = []
        for row in rows:
            if isinstance(row, dict):
                normalized_rows.append([row.get(header, "") for header in headers])
            elif isinstance(row, list):
                normalized_rows.append(row)
            else:
                normalized_rows.append([row])

        return headers, normalized_rows

    lines = _body_lines(body)
    return ["Item", "Details"], [[index, line] for index, line in enumerate(lines, start=1)]


def generate_xlsx_artifact(title: str, body: str, content: dict[str, Any] | None = None) -> Path:
    path = _artifact_path(title or "workbook", "xlsx")
    Workbook = _load_workbook_class()
    workbook = Workbook()
    content = content or {}
    raw_sheets = content.get("sheets")
    sheets = raw_sheets if isinstance(raw_sheets, list) and raw_sheets else [
        {
            "name": "Summary",
            "columns": ["Item", "Details"],
            "rows": [[index, line] for index, line in enumerate(_body_lines(body), start=1)],
        }
    ]

    for index, sheet in enumerate(sheets):
        sheet = sheet if isinstance(sheet, dict) else {}
        ws = workbook.active if index == 0 else workbook.create_sheet()
        ws.title = str(sheet.get("name") or ("Summary" if index == 0 else f"Sheet {index + 1}"))[:31]
        headers, rows = _normalize_sheet_rows(sheet, body)
        ws.append(headers)
        for row in rows:
            ws.append(row)
        ws.freeze_panes = "A2"
        for column_cells in ws.columns:
            max_length = max(len(str(cell.value or "")) for cell in column_cells)
            ws.column_dimensions[column_cells[0].column_letter].width = min(max(max_length + 2, 12), 48)

    workbook.save(path)
    return path


def _fallback_slides(title: str, body: str) -> list[dict[str, Any]]:
    lines = _body_lines(body)
    if not lines:
        lines = ["Generated by Reins WorkMode."]

    slides: list[dict[str, Any]] = [
        {
            "title": title or "WorkMode Presentation",
            "subtitle": "Generated by Reins WorkMode",
            "bullets": [],
        },
        {
            "title": "Overview",
            "bullets": lines[:5],
        },
    ]
    remaining = lines[5:]
    while remaining:
        slides.append({"title": "Details", "bullets": remaining[:6]})
        remaining = remaining[6:]
    return slides


def generate_pptx_artifact(title: str, body: str, content: dict[str, Any] | None = None) -> Path:
    path = _artifact_path(title or "presentation", "pptx")
    Presentation = _load_presentation_class()
    presentation = Presentation()
    content = content or {}
    raw_slides = content.get("slides")
    slides = raw_slides if isinstance(raw_slides, list) and raw_slides else _fallback_slides(title, body)

    for index, slide_data in enumerate(slides):
        slide_data = slide_data if isinstance(slide_data, dict) else {}
        bullets = slide_data.get("bullets")
        has_bullets = isinstance(bullets, list) and len(bullets) > 0
        layout = presentation.slide_layouts[0] if index == 0 and not has_bullets else presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = str(slide_data.get("title") or title or "WorkMode Presentation")

        if index == 0 and not has_bullets and len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = str(slide_data.get("subtitle") or "Generated by Reins WorkMode")
            continue

        if not isinstance(bullets, list):
            bullets = _body_lines(str(slide_data.get("body") or body))[:6]

        if len(slide.placeholders) > 1:
            frame = slide.placeholders[1].text_frame
            frame.clear()
            for bullet_index, bullet in enumerate(bullets):
                paragraph = frame.paragraphs[0] if bullet_index == 0 else frame.add_paragraph()
                paragraph.text = str(bullet)
                paragraph.level = 0

    presentation.save(path)
    return path


def generate_office_artifact(
    artifact_format: str,
    title: str,
    body: str,
    content: dict[str, Any] | None = None,
) -> Path:
    normalized = normalize_office_format(artifact_format) or "docx"

    if normalized == "xlsx":
        return generate_xlsx_artifact(title, body, content)

    if normalized == "pptx":
        return generate_pptx_artifact(title, body, content)

    return generate_docx_artifact(title, body, content)
