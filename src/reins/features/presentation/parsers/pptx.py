from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation

from reins.features.presentation.models import (
    PresentationPlan,
    PresentationRequest,
    PresentationSlide,
    SlideElement,
    SlideElementType,
    SlideLayout,
)


class PresentationIntakeError(ValueError):
    pass


def _clean(value: Any) -> str:
    return " ".join(str(value or "").split()).strip()


def _paragraph_texts(text_frame) -> list[str]:
    values: list[str] = []
    for paragraph in text_frame.paragraphs:
        text = _clean(paragraph.text)
        if text:
            values.append(text)
    return values


def _walk_shapes(shapes):
    for shape in shapes:
        yield shape
        nested = getattr(shape, "shapes", None)
        if nested is not None:
            yield from _walk_shapes(nested)


def extract_pptx_inventory(path: Path) -> dict[str, Any]:
    source = path.expanduser().resolve()
    if not source.is_file() or source.suffix.lower() != ".pptx":
        raise PresentationIntakeError(
            "Presentation editing currently requires a .pptx source file."
        )

    try:
        presentation = Presentation(source)
    except Exception as exc:
        raise PresentationIntakeError(
            f"Could not open the PowerPoint source: {exc}"
        ) from exc

    slides: list[dict[str, Any]] = []
    counts = {
        "text_slots": 0,
        "table_cells": 0,
        "pictures": 0,
        "charts": 0,
    }

    for slide_index, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        title_shape_id = getattr(title_shape, "shape_id", None)
        slots: list[dict[str, Any]] = []

        for shape in _walk_shapes(slide.shapes):
            shape_id = int(getattr(shape, "shape_id", 0) or 0)
            shape_name = str(getattr(shape, "name", "") or "")

            if getattr(shape, "has_table", False):
                table = shape.table
                for row_index, row in enumerate(table.rows):
                    for column_index, cell in enumerate(row.cells):
                        text = "\n".join(_paragraph_texts(cell.text_frame))
                        if not text:
                            continue
                        slots.append(
                            {
                                "target": "table_cell",
                                "shape_id": shape_id,
                                "shape_name": shape_name,
                                "row": row_index,
                                "column": column_index,
                                "role": "table_cell",
                                "text": text,
                            }
                        )
                        counts["table_cells"] += 1
                continue

            if getattr(shape, "has_chart", False):
                counts["charts"] += 1

            if int(getattr(shape, "shape_type", 0) or 0) == 13:
                counts["pictures"] += 1

            if not getattr(shape, "has_text_frame", False):
                continue

            text = "\n".join(_paragraph_texts(shape.text_frame))
            if not text:
                continue

            slots.append(
                {
                    "target": "text_shape",
                    "shape_id": shape_id,
                    "shape_name": shape_name,
                    "role": "title" if shape_id == title_shape_id else "body",
                    "text": text,
                }
            )
            counts["text_slots"] += 1

        slides.append(
            {
                "slide_index": slide_index,
                "slots": slots,
            }
        )

    width = int(presentation.slide_width)
    height = int(presentation.slide_height)
    return {
        "schema": "reins_pptx_inventory.v1",
        "source": str(source),
        "slide_count": len(presentation.slides),
        "aspect_ratio": "4:3" if height and width / height < 1.5 else "16:9",
        "counts": counts,
        "slides": slides,
    }


def extract_pptx_plan(
    path: Path,
    request: PresentationRequest,
) -> PresentationPlan:
    inventory = extract_pptx_inventory(path)
    presentation = Presentation(path)
    slides: list[PresentationSlide] = []

    for source_slide in inventory["slides"]:
        slide_index = int(source_slide["slide_index"])
        slots = list(source_slide["slots"])
        title_slot = next(
            (slot for slot in slots if slot.get("role") == "title"),
            None,
        )
        title = _clean(title_slot.get("text")) if title_slot else ""
        if not title:
            title = next(
                (_clean(slot.get("text")) for slot in slots if slot.get("text")),
                f"Slide {slide_index}",
            )

        body_items: list[str] = []
        for slot in slots:
            if slot is title_slot:
                continue
            for line in str(slot.get("text") or "").splitlines():
                cleaned = _clean(line)
                if cleaned:
                    body_items.append(cleaned)

        notes = ""
        try:
            notes = presentation.slides[slide_index - 1].notes_slide.notes_text_frame.text
        except (AttributeError, ValueError, IndexError):
            pass

        slides.append(
            PresentationSlide(
                index=slide_index,
                type="title" if slide_index == 1 and not body_items else "content",
                title=title,
                elements=[
                    SlideElement(
                        type=SlideElementType.BULLETS,
                        items=body_items,
                    )
                ] if body_items else [],
                speaker_notes=notes,
                layout=SlideLayout(name="source", columns=1),
                metadata={
                    "source_slots": len(slots),
                    "preserve_all_text": True,
                },
            )
        )

    deck_title = _clean(presentation.core_properties.title)
    if not deck_title and slides:
        deck_title = slides[0].title

    return PresentationPlan(
        title=deck_title or path.stem,
        audience=request.audience,
        language=request.language,
        style=request.style,
        aspect_ratio=str(inventory["aspect_ratio"]),
        slides=slides,
        metadata={
            "planner": "pptx_intake",
            "source_pptx": str(path),
            "source_counts": inventory["counts"],
        },
    )
