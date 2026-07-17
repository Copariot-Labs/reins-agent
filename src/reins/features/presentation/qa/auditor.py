from __future__ import annotations

import re
import zipfile

from pathlib import Path
from typing import Any


def audit_presentation_artifact(
    path: Path,
    *,
    expected_slide_count: int,
) -> dict[str, Any]:
    suffix = path.suffix.lower()

    if suffix == ".pptx":
        return _audit_pptx(
            path,
            expected_slide_count=expected_slide_count,
        )

    if suffix in {".html", ".htm"}:
        return _audit_html(
            path,
            expected_slide_count=expected_slide_count,
        )

    if suffix == ".pdf":
        return _audit_pdf(path)

    return {
        "ok": False,
        "checks": [],
        "warnings": [],
        "errors": [f"Unsupported presentation artifact: {path.suffix}"],
    }


def _audit_pptx(
    path: Path,
    *,
    expected_slide_count: int,
) -> dict[str, Any]:
    from pptx import Presentation

    checks: list[dict[str, Any]] = []
    warnings: list[str] = []
    errors: list[str] = []

    if not path.is_file() or path.stat().st_size == 0:
        return {
            "ok": False,
            "checks": checks,
            "warnings": warnings,
            "errors": ["The PPTX output is missing or empty."],
        }

    try:
        with zipfile.ZipFile(path) as archive:
            corrupt_member = archive.testzip()
            if corrupt_member:
                errors.append(
                    f"The PPTX package contains a corrupt member: {corrupt_member}"
                )
            else:
                checks.append({"name": "package_integrity", "passed": True})
    except zipfile.BadZipFile as exc:
        errors.append(f"The PPTX package is invalid: {exc}")

    try:
        presentation = Presentation(path)
    except Exception as exc:
        errors.append(f"PowerPoint could not open the generated deck: {exc}")
        return {
            "ok": False,
            "checks": checks,
            "warnings": warnings,
            "errors": errors,
        }

    actual_count = len(presentation.slides)
    checks.append(
        {
            "name": "slide_count",
            "passed": actual_count == expected_slide_count,
            "expected": expected_slide_count,
            "actual": actual_count,
        }
    )
    if actual_count != expected_slide_count:
        errors.append(
            f"Expected {expected_slide_count} slides, found {actual_count}."
        )

    for index, slide in enumerate(presentation.slides, start=1):
        text_blocks = _collect_shape_text(slide.shapes)
        if not text_blocks:
            warnings.append(f"Slide {index} contains no visible text.")

        for text in text_blocks:
            if len(text) > 700:
                warnings.append(
                    f"Slide {index} has an unusually dense text block ({len(text)} characters)."
                )

    checks.append(
        {
            "name": "non_empty_slides",
            "passed": not any("no visible text" in item for item in warnings),
        }
    )

    return {
        "ok": not errors,
        "checks": checks,
        "warnings": warnings,
        "errors": errors,
        "artifact_size": path.stat().st_size,
    }


def _collect_shape_text(shapes) -> list[str]:
    blocks: list[str] = []

    for shape in shapes:
        if getattr(shape, "has_text_frame", False):
            text = " ".join(shape.text_frame.text.split())
            if text:
                blocks.append(text)

        nested = getattr(shape, "shapes", None)
        if nested is not None:
            blocks.extend(_collect_shape_text(nested))

    return blocks


def _audit_html(
    path: Path,
    *,
    expected_slide_count: int,
) -> dict[str, Any]:
    if not path.is_file() or path.stat().st_size == 0:
        return {
            "ok": False,
            "checks": [],
            "warnings": [],
            "errors": ["The HTML output is missing or empty."],
        }

    source = path.read_text(encoding="utf-8")
    slide_count = len(
        re.findall(r'<section\s+class="slide\s', source)
    )
    errors = []
    if slide_count != expected_slide_count:
        errors.append(
            f"Expected {expected_slide_count} slides, found {slide_count}."
        )
    if "aria-label=" not in source:
        errors.append("The HTML presentation is missing accessible labels.")

    return {
        "ok": not errors,
        "checks": [
            {
                "name": "slide_count",
                "passed": slide_count == expected_slide_count,
                "expected": expected_slide_count,
                "actual": slide_count,
            },
            {
                "name": "keyboard_navigation",
                "passed": "ArrowRight" in source and "ArrowLeft" in source,
            },
        ],
        "warnings": [],
        "errors": errors,
        "artifact_size": path.stat().st_size,
    }


def _audit_pdf(path: Path) -> dict[str, Any]:
    if not path.is_file() or path.stat().st_size < 8:
        return {
            "ok": False,
            "checks": [],
            "warnings": [],
            "errors": ["The PDF output is missing or empty."],
        }

    header = path.read_bytes()[:8]
    valid = header.startswith(b"%PDF-")
    return {
        "ok": valid,
        "checks": [{"name": "pdf_header", "passed": valid}],
        "warnings": [
            "PDF page count is delegated to the LibreOffice conversion result."
        ] if valid else [],
        "errors": [] if valid else ["The converted file is not a valid PDF."],
        "artifact_size": path.stat().st_size,
    }
