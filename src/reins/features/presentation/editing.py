from __future__ import annotations

import hashlib
import json
import re
import shutil
import subprocess
import tempfile
import zipfile

from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from pptx import Presentation

from reins.features.office.content_writer import OfficeContentError, call_reins_json
from reins.features.presentation.engines.frontend_slides import (
    FrontendSlidesEngine,
)
from reins.features.presentation.engines.visuals import get_palette, slugify
from reins.features.presentation.models import (
    PresentationArtifact,
    PresentationEngine,
    PresentationOutputFormat,
    PresentationPlan,
    PresentationRequest,
    PresentationResult,
    PresentationStyle,
)
from reins.features.presentation.parsers import extract_pptx_inventory
from reins.features.presentation.storage import write_json


PML = "http://schemas.openxmlformats.org/presentationml/2006/main"
DML = "http://schemas.openxmlformats.org/drawingml/2006/main"
NS = {"p": PML, "a": DML}


class PresentationEditError(RuntimeError):
    pass


STYLE_FONTS: dict[PresentationStyle, tuple[str, str, str]] = {
    PresentationStyle.MODERN: ("Aptos Display", "Aptos", "Microsoft YaHei"),
    PresentationStyle.TECH: ("Aptos Display", "Aptos", "Microsoft YaHei"),
    PresentationStyle.CORPORATE: ("Arial", "Arial", "Microsoft YaHei"),
    PresentationStyle.CREATIVE: ("Aptos Display", "Aptos", "Microsoft YaHei"),
    PresentationStyle.MINIMAL: ("Arial", "Arial", "Microsoft YaHei"),
    PresentationStyle.DARK: ("Aptos Display", "Aptos", "Microsoft YaHei"),
}


def _sha256(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _package_entries(path: Path) -> dict[str, bytes]:
    try:
        with zipfile.ZipFile(path) as archive:
            if archive.testzip():
                raise PresentationEditError("The PowerPoint package is corrupt.")
            return {
                info.filename: archive.read(info.filename)
                for info in archive.infolist()
                if not info.is_dir()
            }
    except zipfile.BadZipFile as exc:
        raise PresentationEditError("The source is not a valid PPTX package.") from exc


def _write_package(
    source: Path,
    output: Path,
    replacements: dict[str, bytes],
) -> None:
    output.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(source) as source_zip, zipfile.ZipFile(
        output,
        "w",
        compression=zipfile.ZIP_DEFLATED,
    ) as output_zip:
        for info in source_zip.infolist():
            if info.is_dir():
                output_zip.writestr(info, b"")
                continue
            output_zip.writestr(
                info,
                replacements.get(info.filename, source_zip.read(info.filename)),
            )


def _validate_output(path: Path, expected_slides: int) -> None:
    try:
        with zipfile.ZipFile(path) as archive:
            corrupt = archive.testzip()
            if corrupt:
                raise PresentationEditError(
                    f"Edited PPTX contains a corrupt package part: {corrupt}"
                )
        if len(Presentation(path).slides) != expected_slides:
            raise PresentationEditError(
                "Edited PPTX changed the source slide count unexpectedly."
            )
    except (zipfile.BadZipFile, KeyError) as exc:
        raise PresentationEditError("Edited PPTX failed package validation.") from exc


def _preservation_report(
    *,
    source: Path,
    output: Path,
    allowed_parts: set[str],
    operation: str,
) -> dict[str, Any]:
    before = _package_entries(source)
    after = _package_entries(output)
    before_names = set(before)
    after_names = set(after)
    changed = sorted(
        name
        for name in before_names & after_names
        if _sha256(before[name]) != _sha256(after[name])
    )
    added = sorted(after_names - before_names)
    removed = sorted(before_names - after_names)
    unexpected = sorted(set(changed) - allowed_parts)
    ok = not unexpected and not added and not removed
    return {
        "schema": "reins_presentation_preservation.v1",
        "operation": operation,
        "ok": ok,
        "source_sha256": _sha256(source.read_bytes()),
        "output_sha256": _sha256(output.read_bytes()),
        "allowed_changed_parts": sorted(allowed_parts),
        "changed_parts": changed,
        "unexpected_changed_parts": unexpected,
        "added_parts": added,
        "removed_parts": removed,
        "unmodified_parts_preserved": ok,
    }


def _inventory_slots(inventory: dict[str, Any]) -> list[dict[str, Any]]:
    slots: list[dict[str, Any]] = []
    for slide in inventory.get("slides", []):
        for slot in slide.get("slots", []):
            slots.append(
                {
                    **slot,
                    "slide_index": int(slide["slide_index"]),
                }
            )
    return slots


def _deterministic_changes(
    instruction: str,
    slots: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    replacements: list[tuple[str, str]] = []
    patterns = [
        r"replace\s+[\"'“](.+?)[\"'”]\s+with\s+[\"'“](.+?)[\"'”]",
        r"把\s*[\"'“](.+?)[\"'”]\s*(?:改成|改为|替换为)\s*[\"'“](.+?)[\"'”]",
    ]
    for pattern in patterns:
        replacements.extend(re.findall(pattern, instruction, flags=re.IGNORECASE))

    changes: list[dict[str, Any]] = []
    for old, new in replacements:
        for slot in slots:
            text = str(slot.get("text") or "")
            if old not in text:
                continue
            changes.append(
                {
                    **_slot_selector(slot),
                    "old_text": text,
                    "new_text": text.replace(old, new),
                }
            )

    title_patterns = [
        r"(?:change|set)\s+slide\s+(\d+)\s+title\s+(?:to|as)\s+[\"'“]?(.+?)[\"'”]?(?:\.|$)",
        r"第\s*(\d+)\s*页(?:的)?标题\s*(?:改成|改为|替换为)\s*[\"'“]?(.+?)[\"'”]?(?:。|$)",
    ]
    for pattern in title_patterns:
        match = re.search(pattern, instruction, flags=re.IGNORECASE)
        if not match:
            continue
        slide_index = int(match.group(1))
        new_text = match.group(2).strip()
        slot = next(
            (
                value
                for value in slots
                if value["slide_index"] == slide_index
                and value.get("role") == "title"
            ),
            None,
        )
        if slot:
            changes.append(
                {
                    **_slot_selector(slot),
                    "old_text": slot["text"],
                    "new_text": new_text,
                }
            )

    return changes


def _slot_selector(slot: dict[str, Any]) -> dict[str, Any]:
    selector = {
        "slide_index": int(slot["slide_index"]),
        "target": str(slot["target"]),
        "shape_id": int(slot["shape_id"]),
    }
    if slot.get("target") == "table_cell":
        selector.update(
            {
                "row": int(slot["row"]),
                "column": int(slot["column"]),
            }
        )
    return selector


def _model_changes(
    request: PresentationRequest,
    slots: list[dict[str, Any]],
) -> dict[str, Any]:
    inventory_json = json.dumps(
        [
            {
                **_slot_selector(slot),
                "role": slot.get("role"),
                "text": str(slot.get("text") or "")[:800],
            }
            for slot in slots[:400]
        ],
        ensure_ascii=False,
    )
    prompt = f"""
You are planning safe, text-only edits to an existing PowerPoint file.

User instruction:
{request.instruction}

Editable inventory:
{inventory_json}

Return only this JSON shape:
{{
  "summary": "short summary",
  "requires_layout_change": false,
  "reason": "",
  "changes": [
    {{
      "slide_index": 1,
      "target": "text_shape|table_cell",
      "shape_id": 2,
      "row": 0,
      "column": 0,
      "new_text": "complete replacement text for the selected slot"
    }}
  ]
}}

Rules:
- Select only IDs present in the inventory.
- Return the complete new text for each selected slot.
- Do not invent slide or shape IDs.
- Do not alter text unrelated to the instruction.
- If the request changes layout, images, charts, animations, slide order, or slide count,
  set requires_layout_change=true and return no changes.
""".strip()
    try:
        return call_reins_json(
            prompt,
            timeout=int(request.metadata.get("planner_timeout", 180)),
        )
    except OfficeContentError as exc:
        raise PresentationEditError(
            "The edit instruction needs the local presentation model, but it "
            f"was unavailable: {exc}"
        ) from exc


def plan_pptx_modification(
    request: PresentationRequest,
    inventory: dict[str, Any],
) -> dict[str, Any]:
    instruction = str(request.instruction or "").strip()
    slots = _inventory_slots(inventory)
    changes = _deterministic_changes(instruction, slots)
    summary = "Applied explicit text replacements."

    if not changes:
        if request.metadata.get("skip_ai"):
            raise PresentationEditError(
                "The offline editor could not map this instruction to an exact "
                "text target. Use an explicit replacement such as: replace "
                '"old text" with "new text".'
            )
        model_plan = _model_changes(request, slots)
        if model_plan.get("requires_layout_change"):
            reason = str(model_plan.get("reason") or "Layout-aware editing is required.")
            raise PresentationEditError(reason)
        raw_changes = model_plan.get("changes")
        if not isinstance(raw_changes, list):
            raise PresentationEditError("The edit model returned an invalid change plan.")
        changes = [value for value in raw_changes if isinstance(value, dict)]
        summary = str(model_plan.get("summary") or "Updated presentation text.")

    lookup = {
        (
            int(slot["slide_index"]),
            str(slot["target"]),
            int(slot["shape_id"]),
            int(slot.get("row", -1)),
            int(slot.get("column", -1)),
        ): slot
        for slot in slots
    }
    validated: list[dict[str, Any]] = []
    seen: set[tuple[int, str, int, int, int]] = set()

    for change in changes[:100]:
        key = (
            int(change.get("slide_index", 0)),
            str(change.get("target") or "text_shape"),
            int(change.get("shape_id", 0)),
            int(change.get("row", -1)),
            int(change.get("column", -1)),
        )
        slot = lookup.get(key)
        if slot is None:
            raise PresentationEditError(
                f"The edit plan selected an unknown target on slide {key[0]}."
            )
        if key in seen:
            raise PresentationEditError("The edit plan selected the same target twice.")
        seen.add(key)

        old_text = str(slot.get("text") or "")
        new_text = str(change.get("new_text") or "").strip()
        if not new_text:
            raise PresentationEditError(
                "Removing an entire text object requires an explicit deletion workflow."
            )
        if new_text == old_text:
            continue
        old_capacity = max(len(old_text.replace("\n", "")), 12)
        if len(new_text.replace("\n", "")) > old_capacity * 3:
            raise PresentationEditError(
                f"Replacement text on slide {key[0]} is too large for the existing "
                "text frame. Shorten it or use a future layout-aware edit."
            )
        validated.append(
            {
                **_slot_selector(slot),
                "old_text": old_text,
                "new_text": new_text,
            }
        )

    if not validated:
        raise PresentationEditError("The instruction did not produce any text changes.")
    return {
        "schema": "reins_pptx_change_plan.v1",
        "status": "validated",
        "summary": summary,
        "instruction": instruction,
        "changes": validated,
    }


def _shape_container(root: ET.Element, shape_id: int, target: str) -> ET.Element | None:
    tag = "graphicFrame" if target == "table_cell" else "sp"
    for container in root.findall(f".//p:{tag}", NS):
        identity = container.find("./p:nvSpPr/p:cNvPr", NS)
        if identity is None:
            identity = container.find("./p:nvGraphicFramePr/p:cNvPr", NS)
        if identity is not None and int(identity.attrib.get("id", "0")) == shape_id:
            return container
    return None


def _text_paragraphs(container: ET.Element, change: dict[str, Any]) -> list[ET.Element]:
    if change["target"] == "table_cell":
        rows = container.findall(".//a:tbl/a:tr", NS)
        row = int(change["row"])
        column = int(change["column"])
        if row >= len(rows):
            return []
        cells = rows[row].findall("a:tc", NS)
        if column >= len(cells):
            return []
        return cells[column].findall("a:txBody/a:p", NS)
    return container.findall("p:txBody/a:p", NS)


def _paragraph_text(paragraph: ET.Element) -> str:
    return "".join(node.text or "" for node in paragraph.findall(".//a:t", NS))


def _container_text(paragraphs: list[ET.Element]) -> str:
    return "\n".join(_paragraph_text(paragraph) for paragraph in paragraphs).strip()


def _set_paragraph_text(paragraph: ET.Element, value: str) -> None:
    nodes = paragraph.findall(".//a:t", NS)
    if not nodes:
        run = paragraph.find("a:r", NS)
        if run is None:
            run = ET.SubElement(paragraph, f"{{{DML}}}r")
        nodes = [ET.SubElement(run, f"{{{DML}}}t")]
    nodes[0].text = value
    for node in nodes[1:]:
        node.text = ""


def _set_container_text(paragraphs: list[ET.Element], value: str) -> None:
    lines = value.splitlines() or [""]
    if not paragraphs:
        raise PresentationEditError("The selected PowerPoint target has no text body.")
    for index, paragraph in enumerate(paragraphs):
        if index < len(lines):
            line = lines[index]
        else:
            line = ""
        if index == len(paragraphs) - 1 and len(lines) > len(paragraphs):
            line = "\n".join(lines[index:])
        _set_paragraph_text(paragraph, line)


def _blank_target_text(
    root: ET.Element,
    change: dict[str, Any],
) -> None:
    container = _shape_container(
        root,
        int(change["shape_id"]),
        str(change["target"]),
    )
    if container is None:
        raise PresentationEditError("Could not verify the selected text target.")
    for paragraph in _text_paragraphs(container, change):
        for node in paragraph.findall(".//a:t", NS):
            node.text = ""


def modify_pptx(
    *,
    request: PresentationRequest,
    workspace: Path,
) -> PresentationResult:
    assert request.source_path is not None
    source = request.source_path
    inventory = extract_pptx_inventory(source)
    change_plan = plan_pptx_modification(request, inventory)
    change_plan_path = workspace / "change-plan.json"
    write_json(change_plan_path, change_plan)

    replacements: dict[str, bytes] = {}
    changed_parts: set[str] = set()
    changes_by_slide: dict[int, list[dict[str, Any]]] = {}
    for change in change_plan["changes"]:
        changes_by_slide.setdefault(int(change["slide_index"]), []).append(change)

    entries = _package_entries(source)
    for slide_index, changes in changes_by_slide.items():
        part = f"ppt/slides/slide{slide_index}.xml"
        if part not in entries:
            raise PresentationEditError(f"Slide {slide_index} is missing from the PPTX package.")
        root = ET.fromstring(entries[part])
        for change in changes:
            container = _shape_container(
                root,
                int(change["shape_id"]),
                str(change["target"]),
            )
            if container is None:
                raise PresentationEditError(
                    f"Could not find shape {change['shape_id']} on slide {slide_index}."
                )
            paragraphs = _text_paragraphs(container, change)
            current_text = _container_text(paragraphs)
            if current_text != change["old_text"]:
                raise PresentationEditError(
                    f"Slide {slide_index} changed after analysis; the edit was cancelled."
                )
            _set_container_text(paragraphs, str(change["new_text"]))
        verification_source = ET.fromstring(entries[part])
        verification_output = ET.fromstring(ET.tostring(root, encoding="utf-8"))
        for change in changes:
            _blank_target_text(verification_source, change)
            _blank_target_text(verification_output, change)
        if ET.tostring(verification_source) != ET.tostring(verification_output):
            raise PresentationEditError(
                f"Slide {slide_index} changed outside the approved text targets."
            )
        replacements[part] = ET.tostring(root, encoding="utf-8", xml_declaration=True)
        changed_parts.add(part)

    output = workspace / "output" / f"{slugify(source.stem)}-modified.pptx"
    _write_package(source, output, replacements)
    _validate_output(output, int(inventory["slide_count"]))
    report = _preservation_report(
        source=source,
        output=output,
        allowed_parts=changed_parts,
        operation="modify",
    )
    if not report["ok"]:
        output.unlink(missing_ok=True)
        raise PresentationEditError(
            "The edited deck changed package parts outside the approved text targets."
        )
    report["change_count"] = len(change_plan["changes"])
    report["target_scope_verified"] = True
    report_path = workspace / "preservation-report.json"
    write_json(report_path, report)

    return PresentationResult(
        success=True,
        action=request.action,
        engine=PresentationEngine.NATIVE_PPTX,
        primary_output=output,
        artifacts=[
            PresentationArtifact(
                kind="presentation",
                path=output,
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ),
            PresentationArtifact(
                kind="change-plan",
                path=change_plan_path,
                mime_type="application/json",
            ),
            PresentationArtifact(
                kind="preservation-report",
                path=report_path,
                mime_type="application/json",
            ),
        ],
        metadata={"change_count": len(change_plan["changes"])},
    )


def _replace_theme_color(slot: ET.Element, color: str) -> None:
    for child in list(slot):
        slot.remove(child)
    ET.SubElement(slot, f"{{{DML}}}srgbClr", {"val": color})


def _style_from_instruction(value: str | None) -> PresentationStyle | None:
    text = str(value or "").casefold()
    markers = [
        (PresentationStyle.DARK, ("dark", "深色", "黑色")),
        (PresentationStyle.TECH, ("tech", "technology", "科技")),
        (PresentationStyle.CORPORATE, ("corporate", "business", "商务", "企业")),
        (PresentationStyle.CREATIVE, ("creative", "bold", "创意")),
        (PresentationStyle.MINIMAL, ("minimal", "极简")),
        (PresentationStyle.MODERN, ("modern", "现代")),
    ]
    return next(
        (style for style, words in markers if any(word in text for word in words)),
        None,
    )


def restyle_pptx(
    *,
    request: PresentationRequest,
    workspace: Path,
) -> PresentationResult:
    assert request.source_path is not None
    source = request.source_path
    inventory = extract_pptx_inventory(source)
    entries = _package_entries(source)
    resolved_style = _style_from_instruction(request.instruction) or request.style
    palette = get_palette(resolved_style)
    colors = {
        "dk1": palette.text,
        "lt1": palette.background,
        "dk2": palette.primary,
        "lt2": palette.surface,
        "accent1": palette.primary,
        "accent2": palette.accent,
        "accent3": palette.secondary,
        "accent4": palette.muted,
        "accent5": palette.accent,
        "accent6": palette.secondary,
        "hlink": palette.primary,
        "folHlink": palette.accent,
    }
    major_font, minor_font, east_asian_font = STYLE_FONTS[resolved_style]
    replacements: dict[str, bytes] = {}
    allowed_parts: set[str] = set()

    for part, data in entries.items():
        if not part.startswith("ppt/theme/theme") or not part.endswith(".xml"):
            continue
        root = ET.fromstring(data)
        for name, color in colors.items():
            node = root.find(f".//a:clrScheme/a:{name}", NS)
            if node is not None:
                _replace_theme_color(node, color)
        for family, latin_font in (("majorFont", major_font), ("minorFont", minor_font)):
            latin = root.find(f".//a:fontScheme/a:{family}/a:latin", NS)
            east_asian = root.find(f".//a:fontScheme/a:{family}/a:ea", NS)
            complex_script = root.find(f".//a:fontScheme/a:{family}/a:cs", NS)
            if latin is not None:
                latin.set("typeface", latin_font)
            if east_asian is not None:
                east_asian.set("typeface", east_asian_font)
            if complex_script is not None:
                complex_script.set("typeface", latin_font)
        replacements[part] = ET.tostring(root, encoding="utf-8", xml_declaration=True)
        allowed_parts.add(part)

    if not replacements:
        raise PresentationEditError("The source deck does not contain an editable theme part.")

    output = workspace / "output" / f"{slugify(source.stem)}-{resolved_style.value}.pptx"
    _write_package(source, output, replacements)
    _validate_output(output, int(inventory["slide_count"]))
    report = _preservation_report(
        source=source,
        output=output,
        allowed_parts=allowed_parts,
        operation="restyle",
    )
    if not report["ok"]:
        output.unlink(missing_ok=True)
        raise PresentationEditError(
            "The restyled deck changed package parts outside the theme scope."
        )
    report.update(
        {
            "style": resolved_style.value,
            "preserved": [
                "slides",
                "masters and layouts",
                "text and notes",
                "images and media",
                "charts and embedded workbooks",
                "animations and transitions",
            ],
        }
    )
    report_path = workspace / "preservation-report.json"
    write_json(report_path, report)

    return PresentationResult(
        success=True,
        action=request.action,
        engine=PresentationEngine.NATIVE_PPTX,
        primary_output=output,
        artifacts=[
            PresentationArtifact(
                kind="presentation",
                path=output,
                mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            ),
            PresentationArtifact(
                kind="preservation-report",
                path=report_path,
                mime_type="application/json",
            ),
        ],
        warnings=[
            "Theme-linked colors and fonts were restyled. Objects with hard-coded "
            "formatting remain unchanged to preserve slide content and geometry."
        ],
    )


def convert_pptx(
    *,
    request: PresentationRequest,
    plan: PresentationPlan,
    workspace: Path,
) -> PresentationResult:
    assert request.source_path is not None
    source = request.source_path
    if request.output_format == PresentationOutputFormat.HTML:
        result = FrontendSlidesEngine().render(
            request=request,
            plan=plan,
            workspace=workspace,
        )
        warning = (
            "HTML conversion preserves slide order, text, and speaker notes in an "
            "accessible web deck; PowerPoint animations, charts, and exact geometry "
            "are not represented in semantic HTML."
        )
        result.warnings.append(warning)
        report = {
            "schema": "reins_presentation_conversion.v1",
            "operation": "convert",
            "mode": "semantic",
            "source_format": "pptx",
            "output_format": "html",
            "slide_count": len(plan.slides),
            "preserved": ["slide order", "visible text", "speaker notes"],
            "not_preserved": [
                "exact geometry",
                "animations and transitions",
                "native charts and embedded workbooks",
            ],
        }
    elif request.output_format == PresentationOutputFormat.PDF:
        executable = shutil.which("soffice") or shutil.which("libreoffice")
        if not executable:
            raise PresentationEditError(
                "PDF conversion requires LibreOffice (soffice) on the server."
            )
        output_dir = workspace / "output"
        output_dir.mkdir(parents=True, exist_ok=True)
        with tempfile.TemporaryDirectory(prefix="reins-ppt-pdf-") as temp_dir:
            completed = subprocess.run(
                [
                    executable,
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    temp_dir,
                    str(source),
                ],
                capture_output=True,
                text=True,
                timeout=180,
                check=False,
            )
            converted = Path(temp_dir) / f"{source.stem}.pdf"
            if completed.returncode != 0 or not converted.is_file():
                raise PresentationEditError(
                    "LibreOffice could not convert the deck to PDF: "
                    + (completed.stderr.strip() or completed.stdout.strip())
                )
            output = output_dir / f"{slugify(plan.title)}.pdf"
            shutil.copy2(converted, output)
        result = PresentationResult(
            success=True,
            action=request.action,
            engine=PresentationEngine.NATIVE_PPTX,
            primary_output=output,
            artifacts=[
                PresentationArtifact(
                    kind="presentation",
                    path=output,
                    mime_type="application/pdf",
                )
            ],
        )
        report = {
            "schema": "reins_presentation_conversion.v1",
            "operation": "convert",
            "mode": "visual",
            "source_format": "pptx",
            "output_format": "pdf",
            "slide_count": len(plan.slides),
            "preserved": ["rendered appearance", "slide order"],
            "not_preserved": ["editability", "animations and transitions"],
        }
    else:
        raise PresentationEditError(
            "PPTX conversion supports HTML or PDF output."
        )

    report_path = workspace / "preservation-report.json"
    write_json(report_path, report)
    result.artifacts.append(
        PresentationArtifact(
            kind="preservation-report",
            path=report_path,
            mime_type="application/json",
        )
    )
    return result
